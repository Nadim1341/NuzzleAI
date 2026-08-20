import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS_DIR = os.path.abspath(os.path.join(SCRIPT_DIR, '../presentation_screenshots'))
OUTPUT_DIR = os.path.abspath(os.path.join(SCRIPT_DIR, '../presentation'))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, 'Nuzzle_UIUX_Presentation.pptx')
ROOT_OUTPUT_FILE = os.path.abspath(os.path.join(SCRIPT_DIR, '../../Nuzzle_UIUX_Presentation.pptx'))

# Color Definitions
COLOR_BG_DARK = RGBColor(13, 12, 11)       # #0D0C0B
COLOR_BG_SLIDE = RGBColor(22, 20, 19)      # #161413
COLOR_BG_CARD = RGBColor(32, 29, 26)       # #201D1A
COLOR_PRIMARY = RGBColor(244, 145, 92)     # #F4915C Coral
COLOR_AI_PURPLE = RGBColor(139, 92, 246)   # #8B5CF6 Purple
COLOR_AI_INDIGO = RGBColor(99, 102, 241)   # #6366F1 Indigo
COLOR_TEXT_MAIN = RGBColor(250, 247, 242)  # #FAF7F2 White/Cream
COLOR_TEXT_MUTED = RGBColor(163, 150, 141) # #A3968D Grey
COLOR_EMERALD = RGBColor(16, 185, 129)     # #10B981 Green
COLOR_ROSE = RGBColor(244, 63, 94)         # #F43F5E Red

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]

def apply_background(slide, color=COLOR_BG_SLIDE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, slide_num, total_slides=10):
    # Top bar shape
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_BG_DARK
    top_bar.line.color.rgb = RGBColor(51, 45, 40)
    top_bar.line.width = Pt(1)

    # Logo / Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(5.0), Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Nuzzle"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = "Segoe UI"

    # Slide Counter
    txBox2 = slide.shapes.add_textbox(Inches(8.5), Inches(0.12), Inches(4.2), Inches(0.45))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    p2.text = f"Slide {slide_num} / {total_slides} • UI/UX Pitch"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.font.name = "Segoe UI"

def add_title(slide, eyebrow, title):
    # Eyebrow
    tx_eye = slide.shapes.add_textbox(Inches(0.6), Inches(0.85), Inches(10.0), Inches(0.3))
    tf_eye = tx_eye.text_frame
    p_eye = tf_eye.paragraphs[0]
    p_eye.text = eyebrow.upper()
    p_eye.font.size = Pt(10)
    p_eye.font.bold = True
    p_eye.font.color.rgb = COLOR_PRIMARY
    p_eye.font.name = "Segoe UI"

    # Title
    tx_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(10.0), Inches(0.6))
    tf_title = tx_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_MAIN
    p_title.font.name = "Segoe UI"

def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = text

def add_card(slide, x, y, w, h, title, desc, title_color=COLOR_PRIMARY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_BG_CARD
    card.line.color.rgb = RGBColor(51, 45, 40)
    card.line.width = Pt(1)

    tb = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12), w - Inches(0.36), h - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(12.5)
    p1.font.color.rgb = title_color
    p1.font.name = "Segoe UI"
    p1.space_after = Pt(4)

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.font.name = "Segoe UI"

# ==========================================
# SLIDE 1: INTRODUCTION
# ==========================================
slide1 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide1, COLOR_BG_DARK)

# Badge
badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.66), Inches(1.1), Inches(4.0), Inches(0.45))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(42, 28, 20)
badge.line.color.rgb = COLOR_PRIMARY
badge.line.width = Pt(1)
tf_b = badge.text_frame
p_b = tf_b.paragraphs[0]
p_b.text = "NUZZLE  |  SESSION-01: UI/UX PRESENTATION (5 MINS)"
p_b.font.size = Pt(10)
p_b.font.bold = True
p_b.font.color.rgb = COLOR_PRIMARY
p_b.font.name = "Segoe UI"
p_b.alignment = PP_ALIGN.CENTER

# Main Title
tb_m = slide1.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.33), Inches(1.6))
tf_m = tb_m.text_frame
tf_m.word_wrap = True
p_m = tf_m.paragraphs[0]
p_m.text = "Nuzzle — The Next-Gen\nAI-Powered Pet Community"
p_m.font.bold = True
p_m.font.size = Pt(36)
p_m.font.color.rgb = COLOR_TEXT_MAIN
p_m.font.name = "Segoe UI"
p_m.alignment = PP_ALIGN.CENTER

# Subtitle
tb_sub = slide1.shapes.add_textbox(Inches(2.0), Inches(3.6), Inches(9.33), Inches(1.0))
tf_sub = tb_sub.text_frame
tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]
p_sub.text = '"A mobile-first social ecosystem empowering pet parents with dedicated pet profiles, 24/7 AI health triage, and community safety networks."'
p_sub.font.italic = True
p_sub.font.size = Pt(15)
p_sub.font.color.rgb = COLOR_TEXT_MUTED
p_sub.font.name = "Segoe UI"
p_sub.alignment = PP_ALIGN.CENTER

# Team Members Section — 3 side-by-side cards
# Team Lead card (center, slightly larger)
team_members = [
    ("Team Lead", "Abu Zafor", COLOR_PRIMARY),
    ("Team Member", "Nadim Rahman", COLOR_AI_PURPLE),
    ("Team Member", "Rizvi Sarker", COLOR_AI_INDIGO),
]

card_w = Inches(3.5)
card_h = Inches(1.1)
card_y = Inches(4.95)
gap = Inches(0.2)
total_w = card_w * 3 + gap * 2
start_x = (prs.slide_width - total_w) / 2

for i, (role, name, color) in enumerate(team_members):
    cx = start_x + i * (card_w + gap)
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_BG_CARD
    card.line.color.rgb = color
    card.line.width = Pt(1.5)

    tb = slide1.shapes.add_textbox(cx, card_y + Inches(0.12), card_w, Inches(0.38))
    tf = tb.text_frame
    p_role = tf.paragraphs[0]
    p_role.text = role.upper()
    p_role.font.size = Pt(9)
    p_role.font.bold = True
    p_role.font.color.rgb = color
    p_role.font.name = "Segoe UI"
    p_role.alignment = PP_ALIGN.CENTER

    tb2 = slide1.shapes.add_textbox(cx, card_y + Inches(0.48), card_w, Inches(0.46))
    tf2 = tb2.text_frame
    p_name = tf2.paragraphs[0]
    p_name.text = name
    p_name.font.size = Pt(16)
    p_name.font.bold = True
    p_name.font.color.rgb = COLOR_TEXT_MAIN
    p_name.font.name = "Segoe UI"
    p_name.alignment = PP_ALIGN.CENTER

# Duration / Tech pill below team cards
pill = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(6.2), Inches(5.33), Inches(0.6))
pill.fill.solid()
pill.fill.fore_color.rgb = RGBColor(42, 28, 20)
pill.line.color.rgb = RGBColor(51, 45, 40)
pill.line.width = Pt(1)
tf_p = pill.text_frame
p_p = tf_p.paragraphs[0]
p_p.text = "Vue.js 3 + Vite   |   UI/UX Session-01   |   Duration: 5 Mins"
p_p.font.size = Pt(11)
p_p.font.bold = False
p_p.font.color.rgb = COLOR_TEXT_MUTED
p_p.font.name = "Segoe UI"
p_p.alignment = PP_ALIGN.CENTER

add_speaker_notes(slide1, "Good day everyone. Today I'm excited to present Nuzzle—a purpose-built mobile social app designed specifically for pets and pet parents, integrating cutting-edge AI features with clean, human-centered UI/UX design.")

# ==========================================
# SLIDE 2: KEY LEARNINGS
# ==========================================
slide2 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide2)
add_header(slide2, 2)
add_title(slide2, "Section 2 • UI/UX Principles & Insights", "Key Learnings from the UI/UX Session")

points_s2 = [
    ("1. Purpose-Driven Information Architecture", "Eliminated visual clutter common in generic social apps. Every component serves an explicit pet welfare need (Health, Lost Alerts, Playdates)."),
    ("2. Warm, Emotionally Resonant Aesthetic", "Adopted warm paper tokens (#FFFCF9) and soft coral (#F4915C) to evoke warmth, care, and trustworthiness rather than cold tech interfaces."),
    ("3. Thumb-Zone Mobile Ergonomics", "Bottom navigation bar, floating action buttons, and swipeable bottom sheets designed for single-hand mobile interactions."),
    ("4. Privacy & Anonymity by Design", "Separation of Owner identity vs Pet identity, giving users seamless Ghost Mode controls without losing app functionality.")
]

y_pos = Inches(1.85)
for title, desc in points_s2:
    add_card(slide2, Inches(0.6), y_pos, Inches(7.4), Inches(1.15), title, desc)
    y_pos += Inches(1.28)

img_feed = os.path.join(SCREENSHOTS_DIR, '01_home_feed.png')
if os.path.exists(img_feed):
    slide2.shapes.add_picture(img_feed, Inches(9.1), Inches(1.75), height=Inches(5.45))

add_speaker_notes(slide2, "From our UI/UX sessions, my biggest takeaway was that design isn't just decoration—it's empathy. We applied thumb-zone navigation, warm calming color palettes, and privacy-first anonymity.")

# ==========================================
# SLIDE 3: PROBLEM STATEMENT & TARGET USERS
# ==========================================
slide3 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide3)
add_header(slide3, 3)
add_title(slide3, "Section 3 • Problem Statement & Personas", "Problem Statement & Target Users")

points_s3 = [
    ("⚠️ The Problem", "Generic social media (Instagram, TikTok) mixes pet content with human clutter. Pet parents lack a dedicated space for health records, emergency lost alerts, and neighborhood pet compatibility."),
    ("👩‍🦰 Persona 1: The Devoted Pet Parent (Alex)", "Wants a dedicated digital journal for her Golden Retriever, needs 24/7 symptom advice, and loves sharing daily milestones with fellow dog lovers."),
    ("👨‍⚕️ Persona 2: The Community Adopter & Vet", "Needs a trustworthy channel to publish adoptable pets, verify medical passports, and offer slot bookings without scheduling chaos.")
]

y_pos = Inches(1.85)
for title, desc in points_s3:
    add_card(slide3, Inches(0.6), y_pos, Inches(7.4), Inches(1.55), title, desc)
    y_pos += Inches(1.7)

img_prof = os.path.join(SCREENSHOTS_DIR, '13_dual_profile.png')
if os.path.exists(img_prof):
    slide3.shapes.add_picture(img_prof, Inches(9.1), Inches(1.75), height=Inches(5.45))

add_speaker_notes(slide3, "Current platforms treat pets as an afterthought. Pet parents struggle with disorganized medical logs and frantic lost-pet searches. PetSocial provides a dedicated home for both pet identity and care.")

# ==========================================
# SLIDE 4: DESIGN PROCESS & USER FLOW
# ==========================================
slide4 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide4)
add_header(slide4, 4)
add_title(slide4, "Section 4 • Design Process & Approach", "Design Process & User Flows")

points_s4 = [
    ("1. User Research & Empathy Mapping", "Identified top friction points: emergency panic when a pet is lost, vet booking delays, and uncertainty over toxic foods."),
    ("2. Dual-Persona Architecture", "Designed a switchable profile system where the Pet has its own followers and posts, decoupled from the owner's private account."),
    ("3. AI-First Interaction Layer", "Positioned PawAI directly in the bottom navigation for 1-tap access to health vision scanning and symptom triage.")
]

y_pos = Inches(1.85)
for title, desc in points_s4:
    add_card(slide4, Inches(0.6), y_pos, Inches(5.8), Inches(1.55), title, desc)
    y_pos += Inches(1.7)

img_pet = os.path.join(SCREENSHOTS_DIR, '14_pet_persona_profile.png')
img_health = os.path.join(SCREENSHOTS_DIR, '15_health_passport.png')
if os.path.exists(img_pet):
    slide4.shapes.add_picture(img_pet, Inches(6.8), Inches(1.75), height=Inches(5.45))
if os.path.exists(img_health):
    slide4.shapes.add_picture(img_health, Inches(9.9), Inches(1.75), height=Inches(5.45))

add_speaker_notes(slide4, "Our user flows prioritize speed and emotion. Pet parents can switch seamlessly between their human profile and pet passport with one tap, keeping medical records always at hand.")

# ==========================================
# SLIDE 5: PAWAI SUITE
# ==========================================
slide5 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide5)
add_header(slide5, 5)
add_title(slide5, "Section 5 • Prototype Walkthrough — AI Suite", "✨ PawAI Pet Intelligence Suite")

points_s5 = [
    ("🔬 PetScan AI (Vision & Biometrics)", "Neural vision analyzes breed purity (98.4%), Body Condition Score (BCS 5/9), coat health, and mood detection in real-time.", COLOR_AI_PURPLE),
    ("🩺 PawDoctor 24/7 Triage Chat", "Immediate veterinary-trained symptom triage with automated severity flags (Low, Moderate, Critical Emergency).", COLOR_AI_PURPLE),
    ("🎙️ Voice & Emotion AI Translator", "Captures bark/meow audio waveforms and translates pet feelings into entertaining, heartwarming thoughts.", COLOR_AI_PURPLE)
]

y_pos = Inches(1.85)
for title, desc, col in points_s5:
    add_card(slide5, Inches(0.6), y_pos, Inches(5.8), Inches(1.55), title, desc, col)
    y_pos += Inches(1.7)

img_scan = os.path.join(SCREENSHOTS_DIR, '03_pawai_scanner.png')
img_triage = os.path.join(SCREENSHOTS_DIR, '04_pawai_triage.png')
if os.path.exists(img_scan):
    slide5.shapes.add_picture(img_scan, Inches(6.8), Inches(1.75), height=Inches(5.45))
if os.path.exists(img_triage):
    slide5.shapes.add_picture(img_triage, Inches(9.9), Inches(1.75), height=Inches(5.45))

add_speaker_notes(slide5, "Here is our standout innovation: the PawAI suite. Unlike Instagram, we provide active computer vision diagnostics and instant emergency triage when a pet ingests something harmful.")

# ==========================================
# SLIDE 6: SOCIAL AI & PLAYDATES
# ==========================================
slide6 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide6)
add_header(slide6, 6)
add_title(slide6, "Section 5 • Prototype Walkthrough — Social & AI Playdates", "AI Playdate Matcher & Voice Translator")

points_s6 = [
    ("⚡ Playdate Harmony Engine", "Calculates behavioral synergy (e.g. 96% Match between Waffles & Oliver) by matching energy levels and play styles.", COLOR_AI_INDIGO),
    ("💬 Direct Chat with Instant Scheduling", "1-tap playdate invitations, encrypted messaging, and simulated real-time neighborhood coordination.", COLOR_AI_INDIGO),
    ("🎨 Magic Pet Portrait Studio", "Generative AI transforms pet photos into Pixar 3D, Cyberpunk, and Renaissance fine art avatars.", COLOR_AI_INDIGO)
]

y_pos = Inches(1.85)
for title, desc, col in points_s6:
    add_card(slide6, Inches(0.6), y_pos, Inches(5.8), Inches(1.55), title, desc, col)
    y_pos += Inches(1.7)

img_match = os.path.join(SCREENSHOTS_DIR, '06_pawai_matcher.png')
img_trans = os.path.join(SCREENSHOTS_DIR, '05_pawai_translator.png')
if os.path.exists(img_match):
    slide6.shapes.add_picture(img_match, Inches(6.8), Inches(1.75), height=Inches(5.45))
if os.path.exists(img_trans):
    slide6.shapes.add_picture(img_trans, Inches(9.9), Inches(1.75), height=Inches(5.45))

add_speaker_notes(slide6, "Socializing pets can be risky if energy levels clash. Our AI Harmony Matcher checks behavioral compatibility before dogs meet in the park, ensuring safe and fun playdates.")

# ==========================================
# SLIDE 7: CONTENT & STORIES
# ==========================================
slide7 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide7)
add_header(slide7, 7)
add_title(slide7, "Section 5 • Prototype Walkthrough — Media & Feed", "Stories, Reels & AI First-Person Captions")

points_s7 = [
    ("📸 Interactive Story Rings & Viewer", "Auto-advancing 24h stories with progress bars, tap navigation, hold-to-pause, and direct story replies."),
    ("✍️ AI First-Person Caption Writer", "1-tap generative caption tones (Silly, Sweet, Dramatic, Poetic) written from the pet's perspective."),
    ("❤️ Haptic Double-Tap Micro-Interactions", "Double-tap photo for floating heart burst animation, hashtag discovery, and identity-switched comments.")
]

y_pos = Inches(1.85)
for title, desc in points_s7:
    add_card(slide7, Inches(0.6), y_pos, Inches(5.8), Inches(1.55), title, desc)
    y_pos += Inches(1.7)

img_story = os.path.join(SCREENSHOTS_DIR, '02_story_viewer.png')
img_post_ai = os.path.join(SCREENSHOTS_DIR, '16_post_creation_ai.png')
if os.path.exists(img_story):
    slide7.shapes.add_picture(img_story, Inches(6.8), Inches(1.75), height=Inches(5.45))
if os.path.exists(img_post_ai):
    slide7.shapes.add_picture(img_post_ai, Inches(9.9), Inches(1.75), height=Inches(5.45))

add_speaker_notes(slide7, "Content creation is fun and effortless. Our built-in AI Caption Writer generates witty captions in the pet's voice, driving higher engagement across the community.")

# ==========================================
# SLIDE 8: SAFETY & ADOPTION
# ==========================================
slide8 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide8)
add_header(slide8, 8)
add_title(slide8, "Section 5 • Prototype Walkthrough — Safety & Adoption", "Emergency Alert Network & Adoption Hub")

points_s8 = [
    ("🚨 5-Mile Emergency Broadcast", "Instant lost pet alerts pushed to nearby users with location markers, reward tags, and 1-tap call buttons.", COLOR_ROSE),
    ("🏡 Verified Pet Adoption Center", "Filter by species, temperament tags, and vaccination status with direct shelter inquiry messaging.", COLOR_PRIMARY),
    ("🏥 Conflict-Free Vet Booking", "Real-time slot calendar that prevents double bookings and syncs with the Pet Health Passport.", COLOR_EMERALD)
]

y_pos = Inches(1.85)
for title, desc, col in points_s8:
    add_card(slide8, Inches(0.6), y_pos, Inches(5.8), Inches(1.55), title, desc, col)
    y_pos += Inches(1.7)

img_lf = os.path.join(SCREENSHOTS_DIR, '09_lost_and_found.png')
img_comm = os.path.join(SCREENSHOTS_DIR, '08_explore_communities.png')
if os.path.exists(img_lf):
    slide8.shapes.add_picture(img_lf, Inches(6.8), Inches(1.75), height=Inches(5.45))
if os.path.exists(img_comm):
    slide8.shapes.add_picture(img_comm, Inches(9.9), Inches(1.75), height=Inches(5.45))

add_speaker_notes(slide8, "Safety is central to PetSocial. When a pet goes missing, a 5-mile emergency alert is broadcasted immediately, transforming the neighborhood into an active search network.")

# ==========================================
# SLIDE 9: DESIGN SYSTEM & DARK MODE
# ==========================================
slide9 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide9)
add_header(slide9, 9)
add_title(slide9, "Section 5 • Design System & Theming", "Design System, Typography & Dark Mode")

points_s9 = [
    ("🎨 Harmonious Color Palette", "Warm paper background (#FFFCF9), Coral Primary (#F4915C), Emerald for verified health, and Obsidian Dark Mode.", COLOR_EMERALD),
    ("🔤 Modern Typography Scale", "Plus Jakarta Sans for crisp readability + Outfit for bold, friendly headings and brand identity.", COLOR_PRIMARY),
    ("🌙 True Dark Mode (#181615)", "Reduced eye strain during nighttime pet logs and evening walks, preserving battery life on OLED screens.", COLOR_AI_PURPLE)
]

y_pos = Inches(1.85)
for title, desc, col in points_s9:
    add_card(slide9, Inches(0.6), y_pos, Inches(7.4), Inches(1.55), title, desc, col)
    y_pos += Inches(1.7)

img_dark = os.path.join(SCREENSHOTS_DIR, '17_dark_theme_mode.png')
if os.path.exists(img_dark):
    slide9.shapes.add_picture(img_dark, Inches(9.1), Inches(1.75), height=Inches(5.45))

add_speaker_notes(slide9, "We built a robust design token system supporting both warm daylight mode and high-contrast dark mode with zero layout shift.")

# ==========================================
# SLIDE 10: CONCLUSION & FUTURE ROADMAP
# ==========================================
slide10 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide10)
add_header(slide10, 10)
add_title(slide10, "Section 6 • Conclusion & Next Steps", "Conclusion & Future Roadmap")

points_s10 = [
    ("🎓 What We Learned", "Domain-specific social UX succeeds when it solves real emotional and practical problems (health, safety, identity) rather than just copying generic feeds.", COLOR_PRIMARY),
    ("🚀 Planned Next Steps", "• IoT Smart Collar GPS integration for live walk tracking\n• Real-time audio model integration for voice translator\n• Telehealth video calls with licensed veterinary clinics", COLOR_AI_PURPLE),
    ("🐾 Closing Remark", '"Pets are family. PetSocial gives them the dedicated, intelligent platform they deserve."', COLOR_EMERALD)
]

y_pos = Inches(1.85)
for title, desc, col in points_s10:
    add_card(slide10, Inches(0.6), y_pos, Inches(6.8), Inches(1.55), title, desc, col)
    y_pos += Inches(1.7)

# Right Box
ty_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.85), Inches(4.9), Inches(4.8))
ty_box.fill.solid()
ty_box.fill.fore_color.rgb = COLOR_BG_DARK
ty_box.line.color.rgb = COLOR_PRIMARY
ty_box.line.width = Pt(1.5)

tf_ty = ty_box.text_frame
tf_ty.word_wrap = True
p_ty1 = tf_ty.paragraphs[0]
p_ty1.text = "\nThank You! 🐾"
p_ty1.font.bold = True
p_ty1.font.size = Pt(28)
p_ty1.font.color.rgb = COLOR_PRIMARY
p_ty1.font.name = "Segoe UI"
p_ty1.alignment = PP_ALIGN.CENTER

p_ty2 = tf_ty.add_paragraph()
p_ty2.text = "Questions & Live Prototype Demo\n\n"
p_ty2.font.size = Pt(14)
p_ty2.font.color.rgb = COLOR_TEXT_MUTED
p_ty2.font.name = "Segoe UI"
p_ty2.alignment = PP_ALIGN.CENTER

p_ty3 = tf_ty.add_paragraph()
p_ty3.text = "✨ PetSocial Vue 3 Prototype"
p_ty3.font.bold = True
p_ty3.font.size = Pt(13)
p_ty3.font.color.rgb = COLOR_TEXT_MAIN
p_ty3.font.name = "Segoe UI"
p_ty3.alignment = PP_ALIGN.CENTER

add_speaker_notes(slide10, "Thank you for your time. PetSocial demonstrates how UI/UX and AI can combine to create meaningful, joyful products. I welcome any questions and would love to show the live prototype!")

# Save PowerPoint Presentation
os.makedirs(OUTPUT_DIR, exist_ok=True)
prs.save(OUTPUT_FILE)
shutil.copyfile(OUTPUT_FILE, ROOT_OUTPUT_FILE)
print(f"✅ Generated Standard Microsoft PowerPoint Presentation (.pptx):\n1) {OUTPUT_FILE}\n2) {ROOT_OUTPUT_FILE}")
